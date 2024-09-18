from matplotlib.backends.backend_pdf import PdfPages
import matplotlib.pyplot as plt
import os

#plot1 = plotGraph(tempDLstats, tempDLlabels)
#plot2 = plotGraph(tempDLstats_1, tempDLlabels_1)
#plot3 = plotGraph(tempDLstats_2, tempDLlabels_2)

class Reporting:

    def __init__(self, plot):
        
        self.plot = plot
        self.plotPath = plot.fetch.plotFolderPath+plot.fetch.example_folder_path()+'/report'
        os.mkdir(self.plotPath) if not os.path.isdir(self.plotPath) else None #create report folder, in case it does not exist

    def powerpoint(self):
        
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation()
        #blank_slide_layout = prs.slide_layouts[0]
        title_slide_layout = prs.slide_layouts[1]

        top = Inches(1.7)
        center = Inches(3)
        left = Inches(0.3)
        right = Inches(5)
        height = Inches(4.5)

        titles = ['Sanity check', 'Volumes vs returns', 'rho predicts wealth linearly']
        figures = [['Success_parameters', 'Success_parameters'],
        ['Success_parameters'],
        ['Success_parameters', 'Success_parameters']]
        
        positions = [[center], [left, right]]
        
        for slidetitle, figure in zip(titles, figures):
            
            slide = prs.slides.add_slide(title_slide_layout)
            slide.shapes.title.text = slidetitle
            
            position = positions[len(figure)-1]
            
            for n, x in enumerate(position):

                img = self.plot.plotPath+'/'+figure[n]+'.png'            
                pic = slide.shapes.add_picture(img, x, top, height=height)

        prs.save(self.plotPath+'/test.pptx')
        
    def write_report(self):

        plot1 = self.plot.market_dynamics(report=True)
        plot2 = self.plot.attempts_vs_wealth(report=True)
        plot3 = self.plot.success_parameter(report=True)

        pp = PdfPages(self.plotPath+'/training.pdf')
        pp.savefig(plot1)
        pp.savefig(plot2)
        pp.savefig(plot3)
        pp.close()
        
    def merge_reports(self):
        
        from pypdf import PdfMerger

        pdfs = [self.plotPath+'/training.pdf', self.plotPath+'/testing.pdf']

        merger = PdfMerger()

        [merger.append(pdf) for pdf in pdfs]
    
        merger.write("result.pdf")
        merger.close()
